"""
Build PowerPoint presentation for master's thesis defense.
Generates master_work/presentation/presentation.pptx
"""
import io
import os
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = Path(__file__).resolve().parent
IMG_DIR = SCRIPT_DIR / "images"
OUT_PATH = SCRIPT_DIR / "presentation.pptx"

KFU_BLUE = RGBColor(0x00, 0x54, 0x9F)
KFU_DARK = RGBColor(0x00, 0x3D, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
HIGHLIGHT_BG = RGBColor(0xE8, 0xF0, 0xF8)
RED = RGBColor(0xC6, 0x28, 0x28)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def latex_to_image(latex_str, fontsize=18, dpi=200):
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.axis("off")
    fig.patch.set_alpha(0)
    text = ax.text(0, 0, f"${latex_str}$", fontsize=fontsize,
                   ha="left", va="bottom",
                   transform=ax.transAxes)
    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    bbox = text.get_window_extent(renderer)
    w = bbox.width / dpi + 0.15
    h = bbox.height / dpi + 0.15
    fig.set_size_inches(w, h)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, transparent=True,
                bbox_inches="tight", pad_inches=0.03)
    plt.close(fig)
    buf.seek(0)
    return buf


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=18, bold=False, color=BLACK,
             alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(4),
             font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, name="Д.Д. Онищенко", page=""):
    left_box = add_textbox(slide, Inches(0.5), Inches(7.0), Inches(4), Inches(0.4),
                           name, font_size=10, color=GRAY)
    right_box = add_textbox(slide, Inches(9), Inches(7.0), Inches(4), Inches(0.4),
                            page, font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_header(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), fill_color=KFU_BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(12), Inches(0.7),
                title, font_size=28, bold=True, color=WHITE)


def add_subtitle(slide, text):
    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
                text, font_size=16, color=GRAY)


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    path = IMG_DIR / img_name
    if not path.exists():
        return None
    kwargs = {"image_file": str(path), "left": left, "top": top}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    return slide.shapes.add_picture(**kwargs)


def add_formula_image(slide, latex_str, left, top, height=Inches(0.5), fontsize=20):
    buf = latex_to_image(latex_str, fontsize=fontsize)
    pic = slide.shapes.add_picture(buf, left, top, height=height)
    return pic


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=BLACK, bullet_char="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None,
              header=True, highlight_rows=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                         width, Inches(0.35 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if r == 0 and header:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                elif c > 0:
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.alignment = PP_ALIGN.LEFT

            if r == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = KFU_BLUE
            elif highlight_rows and r in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HIGHLIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF8, 0xF8, 0xF8)

    return table_shape


# ============================================================
# SLIDES
# ============================================================

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, KFU_BLUE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.2),
                "Министерство науки и высшего образования РФ\n"
                "Казанский (Приволжский) федеральный университет\n"
                "Институт вычислительной математики и информационных технологий\n"
                "Кафедра прикладной математики и искусственного интеллекта",
                font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.4),
                "Выпускная квалификационная работа магистра",
                font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.6),
                "Обнаружение дефектов на листах картона\n"
                "с использованием моделей семейства YOLO",
                font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Left column: student info
    add_textbox(slide, Inches(1.5), Inches(4.3), Inches(5), Inches(0.3),
                "Студент 2 курса, гр. 09-425", font_size=12, color=RGBColor(0xBB, 0xDE, 0xFB))
    add_textbox(slide, Inches(1.5), Inches(4.6), Inches(5), Inches(0.4),
                "Д.Д. Онищенко", font_size=20, bold=True, color=WHITE)
    add_textbox(slide, Inches(1.5), Inches(5.1), Inches(5), Inches(0.3),
                "Направление", font_size=12, color=RGBColor(0xBB, 0xDE, 0xFB))
    add_textbox(slide, Inches(1.5), Inches(5.4), Inches(5), Inches(0.5),
                "01.04.04 — Прикладная математика\n"
                "Вычислительная геометрия и высокопроизводительные вычисления",
                font_size=13, color=WHITE)

    # Right column: supervisor
    add_textbox(slide, Inches(7.5), Inches(4.3), Inches(5), Inches(0.3),
                "Научный руководитель", font_size=12, color=RGBColor(0xBB, 0xDE, 0xFB))
    add_textbox(slide, Inches(7.5), Inches(4.6), Inches(5), Inches(0.4),
                "А.А. Саламатин", font_size=20, bold=True, color=WHITE)
    add_textbox(slide, Inches(7.5), Inches(5.1), Inches(5), Inches(0.4),
                "к.ф.-м.н., доцент", font_size=14, color=WHITE)

    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
                "Казань · 2026", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "Здравствуйте, уважаемые члены комиссии. Меня зовут Онищенко Даниил, гр. 09-425. "
              "Тема моей работы — «Обнаружение дефектов на листах картона с использованием "
              "моделей семейства YOLO». Научный руководитель — кандидат физико-математических наук, "
              "доцент Саламатин Артур Анатольевич. Расскажу о проблеме, о моём решении и о полученных результатах.")


def slide_02_relevance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Актуальность")
    add_subtitle(slide, "Брак в производстве гофрокартона")

    # Callout boxes
    box1 = add_rect(slide, Inches(0.6), Inches(1.6), Inches(2.8), Inches(1.6),
                    fill_color=RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(2.4), Inches(0.7),
                "10%", font_size=48, bold=True, color=KFU_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(2.4), Inches(0.5),
                "отраслевой стандарт\nдоли брака в ЦБП", font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    box2 = add_rect(slide, Inches(3.6), Inches(1.6), Inches(3.0), Inches(1.6),
                    fill_color=RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide, Inches(3.8), Inches(1.7), Inches(2.6), Inches(0.7),
                "≈ 10 млн ₽", font_size=40, bold=True, color=KFU_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.8), Inches(2.5), Inches(2.6), Inches(0.5),
                "потерь в год\nна типовой завод", font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.6), Inches(3.4), Inches(6.2), Inches(0.4),
                "Расчёт для завода с объёмом 10 млн листов/год при цене ~10 ₽ за лист.",
                font_size=13, color=GRAY)

    # Problem text
    txBox = add_textbox(slide, Inches(0.6), Inches(4.0), Inches(6.2), Inches(2.5), font_size=16)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Проблема:"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = BLACK
    add_para(tf, "Ручной контроль медленный и ошибочный, простые алгоритмы машинного "
             "зрения не справляются с вариативностью текстуры, освещения и форм дефектов.",
             font_size=16, color=BLACK, space_before=Pt(6))
    add_para(tf, "Решение: семейство YOLO — одностадийные детекторы, "
             "способные работать в реальном времени.",
             font_size=16, color=KFU_BLUE, bold=True, space_before=Pt(8))

    # Image on the right
    add_image_safe(slide, "intro.jpg", Inches(7.2), Inches(1.5), width=Inches(5.5))

    add_footer(slide, page="2 / 15")
    set_notes(slide,
              "В целлюлозно-бумажной промышленности доля брака при производстве гофрокартона — около 10% "
              "от общего объёма продукции. Это отраслевой стандарт. Для типового завода с объёмом производства "
              "10 миллионов листов в год при средней цене 10 рублей за лист это даёт около 10 миллионов рублей "
              "потерь ежегодно.\n\nТрадиционные методы — ручной контроль или простые алгоритмы машинного зрения — "
              "не справляются. Поэтому актуальна задача автоматической детекции на основе глубокого обучения. "
              "Среди современных подходов лучше всего подходит семейство YOLO.")


def slide_03_goals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Цель и задачи")

    # Goal box
    goal_box = add_rect(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(1.3),
                        fill_color=KFU_BLUE)
    add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.25),
                "ЦЕЛЬ", font_size=12, color=RGBColor(0xBB, 0xDE, 0xFB))
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.9),
                "Разработать систему автоматической детекции дефектов поверхности "
                "картона в реальном времени на основе легковесных моделей YOLO "
                "и выбрать оптимальную модель для промышленного внедрения",
                font_size=18, color=WHITE)

    add_textbox(slide, Inches(0.6), Inches(2.8), Inches(5), Inches(0.4),
                "ЗАДАЧИ", font_size=14, bold=True, color=GRAY)

    tasks = [
        "1.  Проанализировать архитектуру моделей семейства YOLO и обосновать выбор легковесных версий (nano / tiny)",
        "2.  Подготовить и предобработать датасет (1056 изображений, разбиение 7:2:1)",
        "3.  Выполнить байесовскую оптимизацию гиперпараметров",
        "4.  Сравнить 5 моделей (YOLOv8n, v8n tuned, v9t, v10n, v11n) со статистическим обоснованием выбора",
        "5.  Разработать рекомендации по внедрению на конвейерных линиях и оценить экономический потенциал",
    ]
    txBox = add_textbox(slide, Inches(0.6), Inches(3.2), Inches(12), Inches(3.5), font_size=17)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, task in enumerate(tasks):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = task
        p.font.size = Pt(17)
        p.font.color.rgb = BLACK
        p.font.name = "Calibri"
        p.space_after = Pt(8)

    add_footer(slide, page="3 / 15")
    set_notes(slide,
              "Цель работы — разработать систему автоматической детекции дефектов поверхности картона "
              "в реальном времени на основе легковесных моделей YOLO и выбрать оптимальную модель для "
              "промышленного внедрения.\n\nДля достижения цели поставлено пять задач:\n"
              "первая — проанализировать архитектуру семейства YOLO;\n"
              "вторая — подготовить датасет из 1056 изображений;\n"
              "третья — выполнить байесовскую оптимизацию гиперпараметров;\n"
              "четвёртая — сравнить пять моделей со статистическим обоснованием;\n"
              "пятая — разработать рекомендации по внедрению и оценить экономику.")


def slide_04_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Архитектура YOLOv8")
    add_subtitle(slide, "Backbone → Neck → Head")

    # Architecture image
    add_image_safe(slide, "image30.png", Inches(0.4), Inches(1.6), width=Inches(7.5))
    add_textbox(slide, Inches(0.4), Inches(6.3), Inches(7.5), Inches(0.4),
                "Рисунок 1 — Архитектура YOLOv8", font_size=11, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # Cards on the right
    cards = [
        ("Backbone", [
            "Извлечение признаков",
            "C2f-блоки (Cross Stage Partial)",
            "SPPF — расширение поля зрения",
        ]),
        ("Neck", [
            "FPN (top-down) + PAN (bottom-up)",
            "Объединение признаков разных масштабов",
        ]),
        ("Head", [
            "Anchor-free предсказание",
            "Bbox + класс + уверенность",
        ]),
    ]
    y = Inches(1.5)
    for title, items in cards:
        card = add_rect(slide, Inches(8.2), y, Inches(4.5), Inches(1.5),
                        fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
        add_textbox(slide, Inches(8.4), y + Emu(Inches(0.1).emu), Inches(4), Inches(0.35),
                    title, font_size=16, bold=True, color=KFU_BLUE)
        bullet_y = y + Emu(Inches(0.45).emu)
        for item in items:
            add_textbox(slide, Inches(8.4), bullet_y, Inches(4), Inches(0.3),
                        f"• {item}", font_size=13, color=BLACK)
            bullet_y += Emu(Inches(0.28).emu)
        y += Emu(Inches(1.65).emu)

    add_footer(slide, page="4 / 15")
    set_notes(slide,
              "YOLOv8 — одностадийный детектор, который одновременно выполняет классификацию и "
              "локализацию объектов. Архитектурно состоит из трёх частей.\n\n"
              "Backbone извлекает признаки. Это последовательность свёрток и C2f-блоков. "
              "На выходе backbone стоит SPPF — он расширяет рецептивное поле.\n\n"
              "Neck — это связка FPN и PAN. FPN строит иерархию признаков сверху вниз, "
              "PAN добавляет обратный путь снизу вверх. Это улучшает локализацию мелких объектов.\n\n"
              "Head — anchor-free, выдаёт координаты боксов, класс и уверенность.")


def slide_05_loss_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Функция потерь YOLOv8 · 1 из 2")
    add_subtitle(slide, "Классификация и регрессия боксов")

    # Classification loss block
    block1 = add_rect(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(2.2),
                      fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
                "Классификация — бинарная кросс-энтропия", font_size=14, bold=True, color=KFU_BLUE)
    add_formula_image(slide,
                      r"L_{\mathrm{cls}} = - \left[ y \log(p) + (1 - y) \log(1 - p) \right]",
                      Inches(2.5), Inches(2.1), height=Inches(0.5), fontsize=22)
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.7),
                "p — предсказанная вероятность класса, y ∈ {0, 1} — истинная метка. "
                "Штраф растёт, когда модель ошибается. Для C = 1 (один класс «defect»).",
                font_size=14, color=GRAY)

    # CIoU loss block
    block2 = add_rect(slide, Inches(0.6), Inches(4.1), Inches(12), Inches(2.6),
                      fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.3),
                "Регрессия боксов — Complete IoU", font_size=14, bold=True, color=KFU_BLUE)
    add_formula_image(slide,
                      r"L_{\mathrm{CIoU}} = 1 - \mathrm{IoU} + \frac{\rho^2(b, b^{gt})}{c^2} + \alpha v",
                      Inches(3.0), Inches(4.6), height=Inches(0.6), fontsize=22)
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.0),
                "Учитывает: перекрытие (IoU), расстояние между центрами (ρ²/c²) "
                "и согласованность соотношения сторон (αv). "
                "Улучшает локализацию дефектов по сравнению с обычным IoU.",
                font_size=14, color=GRAY)

    add_footer(slide, page="5 / 15")
    set_notes(slide,
              "При обучении оптимизируется комбинированная функция потерь из трёх компонентов.\n\n"
              "Первый — классификационная потеря на основе бинарной кросс-энтропии. "
              "Она измеряет, насколько правильно модель определяет класс объекта.\n\n"
              "Второй — потеря регрессии боксов, Complete IoU. Помимо площади пересечения "
              "учитывается расстояние между центрами рамок и согласованность соотношения сторон.")


def slide_06_loss_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Функция потерь YOLOv8 · 2 из 2")
    add_subtitle(slide, "Distribution Focal Loss и итоговая комбинация")

    # DFL block
    block1 = add_rect(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(2.2),
                      fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
                "Distribution Focal Loss — позиции как распределение",
                font_size=14, bold=True, color=KFU_BLUE)
    add_formula_image(slide,
                      r"L_{\mathrm{DFL}} = - \left( (y_{i+1} - y) \log(S_i) + (y - y_i) \log(S_{i+1}) \right)",
                      Inches(1.8), Inches(2.15), height=Inches(0.5), fontsize=22)
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.8),
                "Координаты бокса моделируются как дискретное распределение на сетке точек. "
                "Sᵢ — предсказанные вероятности, y лежит между yᵢ и yᵢ₊₁. "
                "Позволяет учитывать неопределённость границ дефектов.",
                font_size=14, color=GRAY)

    # Total loss callout
    total_box = add_rect(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(2.0),
                         fill_color=RGBColor(0xE3, 0xF2, 0xFD),
                         line_color=KFU_BLUE)
    add_textbox(slide, Inches(1.7), Inches(4.35), Inches(9.8), Inches(0.3),
                "ОБЩАЯ ФУНКЦИЯ ПОТЕРЬ", font_size=12, bold=True, color=KFU_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_formula_image(slide,
                      r"L = 7.5 \cdot L_{\mathrm{CIoU}} + 0.5 \cdot L_{\mathrm{cls}} + 1.5 \cdot L_{\mathrm{DFL}}",
                      Inches(3.0), Inches(4.7), height=Inches(0.45), fontsize=22)
    add_textbox(slide, Inches(1.7), Inches(5.4), Inches(9.8), Inches(0.6),
                "Регрессия доминирует (вес 7.5) — точная локализация дефекта критична "
                "для промышленного применения.",
                font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, page="6 / 15")
    set_notes(slide,
              "Третий компонент — Distribution Focal Loss. Координаты бокса моделируются не точечно, "
              "а как дискретное распределение. Это позволяет учитывать неопределённость.\n\n"
              "Общая функция — взвешенная сумма с коэффициентами 7.5 (CIoU), 0.5 (cls) и 1.5 (DFL). "
              "Регрессия имеет наибольший вес, потому что точная локализация дефекта критична для производства.")


def slide_07_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Датасет")
    add_subtitle(slide, "1056 изображений · 640×640 · 1 класс «defect» · split 7:2:1")

    # Three images
    img_x = [Inches(0.6), Inches(4.6), Inches(8.6)]
    img_names = ["image31.png", "image32.png", "image34.png"]
    captions = ["(а) Царапины", "(б) Вмятины", "(в) Разрывы"]
    for x, name, cap in zip(img_x, img_names, captions):
        add_image_safe(slide, name, x, Inches(1.6), width=Inches(3.6))
        add_textbox(slide, x, Inches(5.2), Inches(3.6), Inches(0.3),
                    cap, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(1.0),
                "Pascal VOC → YOLO format. Все типы дефектов объединены в один класс — "
                "для производства важен сам факт дефекта, а не его природа. "
                "Расширенная валидация (20%) — для надёжного подбора гиперпараметров.",
                font_size=16, color=GRAY, bold=False)

    add_footer(slide, page="7 / 15")
    set_notes(slide,
              "Использован открытый датасет — 1056 изображений с разрешением 640 на 640 пикселей. "
              "Три типа дефектов: царапины, вмятины, разрывы.\n\n"
              "Все дефекты объединены в один класс «defect». "
              "Разметка из формата Pascal VOC сконвертирована в YOLO-формат.\n\n"
              "Разбиение 7:2:1 вместо стандартного 8:1:1: увеличенная валидационная выборка 20%.")


def slide_08_hyperparams(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Подбор гиперпараметров")
    add_subtitle(slide, "Байесовская оптимизация")

    # Left: description
    card = add_rect(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(2.2),
                    fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.35),
                "Метод Байеса", font_size=16, bold=True, color=KFU_BLUE)
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5), Inches(1.5), [
        "Гауссовский процесс моделирует зависимость метрики от гиперпараметров",
        "Максимизация ожидаемого улучшения",
        "Эффективнее grid / random search",
    ], font_size=14)

    add_textbox(slide, Inches(0.6), Inches(4.0), Inches(5.5), Inches(0.8),
                "Базовая модель: YOLOv8n — оптимальный баланс скорости и качества; "
                "более крупные версии (s/m/l/x) не дают прироста на малом датасете.",
                font_size=14, color=BLACK)
    add_textbox(slide, Inches(0.6), Inches(4.9), Inches(5.5), Inches(0.5),
                "Финальное обучение: 50 эпох — достаточно для стабильной сходимости.",
                font_size=14, bold=True, color=BLACK)

    # Right: table
    table_data = [
        ["Гиперпараметр", "Диапазон", "Итог"],
        ["lr0", "[10⁻⁵, 10⁻²]", "0.0026"],
        ["lrf", "[0.01, 0.2]", "0.102"],
        ["momentum", "[0.85, 0.95]", "0.939"],
        ["weight decay", "[0, 0.01]", "0.0006"],
        ["warmup epochs", "[1, 3]", "2"],
        ["mosaic", "[0, 1]", "1.0"],
        ["mixup", "[0, 0.3]", "0.025"],
        ["fliplr", "[0, 0.6]", "0.283"],
        ["cos lr", "{F, T}", "True"],
    ]
    add_table(slide, Inches(6.5), Inches(1.5), Inches(6.2), table_data,
              col_widths=[Inches(2.2), Inches(2.0), Inches(2.0)])
    add_textbox(slide, Inches(6.5), Inches(5.8), Inches(6.2), Inches(0.3),
                "Таблица 1 — Ключевые гиперпараметры (всего подбиралось 19)",
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, page="8 / 15")
    set_notes(slide,
              "Для подбора гиперпараметров применён байесовский метод. Он строит вероятностную модель "
              "и выбирает наиболее перспективную комбинацию. Подбирались 19 гиперпараметров. "
              "Финальное обучение — 50 эпох.")


def slide_09_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Сравнение моделей")
    add_subtitle(slide, "Первичные метрики на тестовой выборке")

    # Chart image
    add_image_safe(slide, "image29.png", Inches(0.4), Inches(1.5), width=Inches(6.5))
    add_textbox(slide, Inches(0.4), Inches(6.0), Inches(6.5), Inches(0.3),
                "Рисунок 2 — Метрики precision, recall, mAP50, mAP50-95, FPS",
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Table
    table_data = [
        ["Модель", "mAP50", "FPS"],
        ["YOLOv8n", "0.944", "116"],
        ["YOLOv8n tuned", "0.943", "116"],
        ["YOLOv9t", "0.939", "89"],
        ["YOLOv10n", "0.912", "111"],
        ["YOLOv11n", "0.943", "85"],
    ]
    add_table(slide, Inches(7.3), Inches(1.5), Inches(5.5), table_data,
              col_widths=[Inches(2.5), Inches(1.5), Inches(1.5)],
              highlight_rows=[1])

    add_textbox(slide, Inches(7.3), Inches(4.1), Inches(5.5), Inches(1.5),
                "Все модели близки (precision > 0.89, recall > 0.85). "
                "По точечным замерам выбор неочевиден — нужна "
                "статистическая оценка с учётом вариативности.",
                font_size=15, color=BLACK)

    add_footer(slide, page="9 / 15")
    set_notes(slide,
              "Обучил пять моделей. По первичным замерам все показывают высокое качество. "
              "По FPS лидируют YOLOv8n со 116 кадрами в секунду. "
              "На единичных замерах нельзя надёжно выбрать лучшую модель — "
              "применил статистический анализ с бутстрэпом.")


def slide_10_statistics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Статистический анализ")
    add_subtitle(slide, "Бутстрэп → 4 теста → выбор модели")

    # Left column - tests
    cards_left = [
        ("Bootstrap (B = 20)",
         "Оценка вариативности метрики mAP50 на ограниченной выборке через подвыборки с возвращением."),
        ("Шапиро–Уилка",
         "Проверка нормальности распределения mAP50.\n"
         "v8n p=0.34, v8n tuned p=0.11, v9t p=0.19 — нормальные\n"
         "v10n p=0.019, v11n p=0.016 — отклонение → непараметрический тест"),
        ("Левена",
         "Равенство дисперсий: p = 0.840 — дисперсии равны."),
    ]
    y = Inches(1.5)
    for title, desc in cards_left:
        card = add_rect(slide, Inches(0.5), y, Inches(6.0), Inches(1.35),
                        fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
        add_textbox(slide, Inches(0.7), y + Emu(Inches(0.05).emu), Inches(5.5), Inches(0.3),
                    title, font_size=13, bold=True, color=KFU_BLUE)
        add_textbox(slide, Inches(0.7), y + Emu(Inches(0.35).emu), Inches(5.5), Inches(0.9),
                    desc, font_size=12, color=BLACK)
        y += Emu(Inches(1.5).emu)

    # Right column
    y_r = Inches(1.5)
    card_kw = add_rect(slide, Inches(6.8), y_r, Inches(6.0), Inches(1.1),
                       fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, Inches(7.0), y_r + Emu(Inches(0.05).emu), Inches(5.5), Inches(0.3),
                "Краскела–Уоллиса", font_size=13, bold=True, color=KFU_BLUE)
    add_textbox(slide, Inches(7.0), y_r + Emu(Inches(0.35).emu), Inches(5.5), Inches(0.6),
                "Непараметрический аналог ANOVA.\nH = 46.43, p = 2·10⁻⁹ — различия значимы.",
                font_size=12, color=BLACK)

    # Mann-Whitney table
    y_mw = Inches(2.8)
    add_textbox(slide, Inches(7.0), y_mw, Inches(5.5), Inches(0.3),
                "Манна–Уитни (попарный)", font_size=13, bold=True, color=KFU_BLUE)

    mw_data = [
        ["Пара", "p-значение", "Значимо"],
        ["v10n vs v8n", "6.8·10⁻⁸", "Да"],
        ["v10n vs v8n tuned", "1.7·10⁻⁷", "Да"],
        ["v10n vs v9t", "6.8·10⁻⁸", "Да"],
        ["v10n vs v11n", "3.4·10⁻⁷", "Да"],
        ["v8n vs v8n tuned", "0.69", "—"],
        ["v8n vs v9t", "0.90", "—"],
        ["v8n vs v11n", "0.58", "—"],
    ]
    add_table(slide, Inches(6.8), Inches(3.15), Inches(6.0), mw_data,
              col_widths=[Inches(2.5), Inches(1.8), Inches(1.7)])

    add_footer(slide, page="10 / 15")
    set_notes(slide,
              "Чтобы корректно сравнить модели, применил бутстрэп — 20 подвыборок с возвращением. "
              "Шапиро–Уилка: для v10n и v11n нормальность отвергается → непараметрические тесты. "
              "Левена: дисперсии равны. Краскела–Уоллиса: p = 2·10⁻⁹, различия значимы. "
              "Манна–Уитни: v10n значимо хуже всех; остальные статистически эквивалентны. "
              "Вывод: YOLOv8n с дефолтами — оптимальный выбор (максимальный FPS при эквивалентном качестве).")


def slide_11_final_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Итоговая модель: YOLOv8n")
    add_subtitle(slide, "Графики обучения")

    # Callout boxes
    metrics = [("0.944", "mAP@50"), ("0.682", "mAP@50-95"), ("116", "FPS"), ("50", "эпох")]
    x = Inches(0.6)
    for val, label in metrics:
        box = add_rect(slide, x, Inches(1.5), Inches(2.8), Inches(1.1),
                       fill_color=RGBColor(0xE3, 0xF2, 0xFD))
        add_textbox(slide, x + Inches(0.1), Inches(1.55), Inches(2.6), Inches(0.55),
                    val, font_size=36, bold=True, color=KFU_BLUE, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(2.15), Inches(2.6), Inches(0.3),
                    label, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
        x += Inches(3.0)

    # Training curves
    add_image_safe(slide, "image33.png", Inches(0.3), Inches(2.9), width=Inches(12.5))

    add_footer(slide, page="11 / 15")
    set_notes(slide,
              "Итоговая модель — YOLOv8n с дефолтными гиперпараметрами. "
              "mAP50 = 0.944, mAP50-95 = 0.682 при скорости 116 FPS. "
              "Лоссы устойчиво снижаются, метрики выходят на плато к 30-й эпохе. "
              "Для сравнения: бейзлайн YOLOv5-S-G-B достигает того же mAP50, но за 300 эпох.")


def slide_12_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Рекомендации по внедрению")
    add_subtitle(slide, "Поэтапный подход")

    # Stage 1
    stage1 = add_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2),
                      fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.35),
                "Этап 1 · Односторонний осмотр", font_size=16, bold=True, color=KFU_BLUE)
    items1 = [
        "Камеры: 2–4 промышленные (Basler/IDS, 4K+, 60+ FPS) над конвейером",
        "Подсветка: LED структурированная",
        "Edge: NVIDIA Jetson Orin Nano + TensorRT",
        "Управление: PLC + пневматический отбраковщик",
        "Задержка: 20–50 мс на кадр",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.5),
                    items1, font_size=14)

    result1 = add_rect(slide, Inches(0.8), Inches(4.8), Inches(5.4), Inches(0.7),
                       fill_color=RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide, Inches(1.0), Inches(4.85), Inches(5.0), Inches(0.5),
                "Снижение поверхностного брака с 10% → 1–2%",
                font_size=15, bold=True, color=KFU_BLUE, alignment=PP_ALIGN.CENTER)

    # Stage 2
    stage2 = add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2),
                      fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.35),
                "Этап 2 · Двусторонний осмотр", font_size=16, bold=True, color=KFU_BLUE)
    items2 = [
        "Одновременный: камеры сверху и снизу + Jetson AGX Orin; задержка 30–70 мс",
        "С переворотом: робо-флиппер (Festo/ABB); задержка ~1.5–2.5 с/лист",
        "Детектируется внутренний брак (расслоение, склейка)",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(2.5),
                    items2, font_size=14)

    add_textbox(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(0.8),
                "Переход на этап 2 оправдан, если оборотный брак > 15–20% от общего.",
                font_size=13, color=GRAY)

    add_footer(slide, page="12 / 15")
    set_notes(slide,
              "Рекомендую поэтапное внедрение.\n\n"
              "Первый этап — односторонний осмотр. 2–4 камеры, Jetson Orin + TensorRT. "
              "Снижает поверхностный брак с 10% до 1–2%.\n\n"
              "Второй этап — двусторонний осмотр. Два варианта: одновременный или с переворотом. "
              "Переход оправдан, если доля оборотного брака > 15–20%.")


def slide_13_conveyor(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Схема станции контроля")
    add_subtitle(slide, "Визуализация этапа 1: камеры · edge-AI · отбраковщик")

    add_image_safe(slide, "conveyor.jpg", Inches(1.5), Inches(1.7), width=Inches(10.3))

    add_footer(slide, page="13 / 15")
    set_notes(slide,
              "Над конвейером — промышленные камеры с LED-подсветкой. "
              "Edge-AI модуль Jetson Orin с TensorRT. "
              "Пневматический отбраковщик по сигналу PLC. "
              "Всё работает локально, задержка 20–50 мс на кадр.")


def slide_14_economics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Экономика и заключение")

    # Left: formulas
    formulas = [
        ("Инвестиции",
         r"I(N) = N \cdot C_{hw} + C_{sw}^{(1)} + (N-1) \cdot C_{sw}^{(+)}",
         "N — число линий, C_hw — аппаратура/линия, "
         "C_sw¹ — ПО для 1-й линии, C_sw⁺ — ПО для каждой следующей"),
        ("Чистая годовая экономия / линия",
         r"S_{net} = S_{la} + Q \cdot d_1 \cdot p \cdot \eta - Q(1-d) \cdot p \cdot \alpha - C_{op}",
         "Q — листов/год, d — доля брака, p — цена листа, "
         "η — recall, α — ложные срабатывания"),
        ("Окупаемость и ROI",
         r"T(N) = \frac{12 \cdot I(N)}{N \cdot S_{net}}, \quad ROI(N) = \frac{N \cdot S_{net} - I(N)}{I(N)} \times 100\%",
         "T — месяцы, ROI — возврат инвестиций за 1-й год"),
    ]
    y = Inches(1.15)
    for title, latex, desc in formulas:
        block = add_rect(slide, Inches(0.4), y, Inches(6.6), Inches(1.7),
                         fill_color=LIGHT_GRAY, line_color=RGBColor(0xDD, 0xDD, 0xDD))
        add_textbox(slide, Inches(0.55), y + Emu(Inches(0.05).emu), Inches(6.3), Inches(0.25),
                    title, font_size=12, bold=True, color=KFU_BLUE)
        add_formula_image(slide, latex,
                          Inches(0.7), y + Emu(Inches(0.35).emu),
                          height=Inches(0.5), fontsize=18)
        add_textbox(slide, Inches(0.55), y + Emu(Inches(1.0).emu), Inches(6.3), Inches(0.6),
                    desc, font_size=10, color=GRAY)
        y += Emu(Inches(1.85).emu)

    # Right: table
    econ_data = [
        ["Показатель", "1 линия", "5 линий", "10 линий"],
        ["Инвестиции, млн ₽", "1.7–3.7", "6.9–14.9", "13.4–28.9"],
        ["S_net, млн ₽/год", "5.4–9.1", "27–45", "54–91"],
        ["Окупаемость", "2–8 мес", "2–7 мес", "2–6 мес"],
    ]
    add_table(slide, Inches(7.3), Inches(1.15), Inches(5.5), econ_data,
              col_widths=[Inches(1.8), Inches(1.2), Inches(1.2), Inches(1.3)],
              highlight_rows=[3])
    add_textbox(slide, Inches(7.3), Inches(3.0), Inches(5.5), Inches(0.3),
                "Q = 10 млн листов/год · d = 10% · p = 10 ₽",
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Conclusion grid
    conclusions = [
        ("✓", "Все 5 задач выполнены"),
        ("YOLOv8n", "mAP50 0.944, 116 FPS"),
        ("2–8 мес", "окупаемость"),
        ("→", "300 эпох · больше данных · реальная линия"),
    ]
    cx = Inches(7.3)
    cy = Inches(3.5)
    for i, (num, text) in enumerate(conclusions):
        col = i % 2
        row = i // 2
        bx = cx + Inches(col * 2.8)
        by = cy + Inches(row * 1.6)
        box = add_rect(slide, bx, by, Inches(2.6), Inches(1.4),
                       fill_color=RGBColor(0xE3, 0xF2, 0xFD))
        add_textbox(slide, bx + Inches(0.1), by + Inches(0.05), Inches(2.4), Inches(0.55),
                    num, font_size=28, bold=True, color=KFU_BLUE, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, bx + Inches(0.1), by + Inches(0.7), Inches(2.4), Inches(0.55),
                    text, font_size=13, color=BLACK, alignment=PP_ALIGN.CENTER)

    add_footer(slide, page="14 / 15")
    set_notes(slide,
              "Разработана общая экономическая модель. "
              "Для типовых параметров окупаемость на одну линию — от 2 до 8 месяцев. "
              "При масштабировании на 10 линий — 2–6 месяцев.\n\n"
              "Итог: все пять задач выполнены. YOLOv8n — 0.944 mAP50, 116 FPS. "
              "Перспективы: увеличение эпох до 300, расширение датасета, "
              "тестирование в реальных производственных условиях.\n\n"
              "Спасибо за внимание, готов ответить на ваши вопросы.")


def slide_15_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KFU_BLUE)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
                "Спасибо за внимание!",
                font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.5),
                "Готов ответить на ваши вопросы",
                font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(1.2),
                "Д.Д. Онищенко · гр. 09-425\n"
                "Научный руководитель: к.ф.-м.н., доцент А.А. Саламатин\n"
                "Казань · 2026",
                font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    set_notes(slide, "Финальный слайд. Поблагодарить комиссию, обозначить готовность к вопросам.")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_relevance(prs)
    slide_03_goals(prs)
    slide_04_architecture(prs)
    slide_05_loss_1(prs)
    slide_06_loss_2(prs)
    slide_07_dataset(prs)
    slide_08_hyperparams(prs)
    slide_09_comparison(prs)
    slide_10_statistics(prs)
    slide_11_final_model(prs)
    slide_12_deployment(prs)
    slide_13_conveyor(prs)
    slide_14_economics(prs)
    slide_15_thanks(prs)

    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
